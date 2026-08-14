"""Safe local attachment storage and lightweight content extraction.

Attachments are runtime inputs, never source-of-truth character data.  The
manager keeps files below ``runtime/chat/attachments`` and stores only a
bounded extracted representation in the Agent metadata database.
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone
from hashlib import sha256
import json
import mimetypes
from pathlib import Path
import re
from typing import Any, BinaryIO
from uuid import uuid4

from .agent_store import AgentStore


MAX_FILE_BYTES = 50 * 1024 * 1024
MAX_IMAGE_BYTES = 20 * 1024 * 1024
MAX_AUDIO_BYTES = 500 * 1024 * 1024  # duration is checked when metadata is available
MAX_ATTACHMENTS_PER_TURN = 10
MAX_EXTRACTED_CHARS = 1_000_000

ALLOWED_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".webp", ".gif",
    ".pdf", ".txt", ".md", ".csv", ".json",
    ".docx", ".xlsx", ".pptx",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css",
    ".xml", ".yaml", ".yml", ".toml", ".ini", ".log",
    ".wav", ".mp3", ".m4a", ".webm", ".ogg",
}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".gif"}
AUDIO_EXTENSIONS = {".wav", ".mp3", ".m4a", ".webm", ".ogg"}
DOCUMENT_EXTENSIONS = ALLOWED_EXTENSIONS - IMAGE_EXTENSIONS - AUDIO_EXTENSIONS


class AttachmentError(ValueError):
    """A user-correctable attachment validation or parsing error."""


@dataclass(frozen=True)
class AttachmentLimits:
    max_file_bytes: int = MAX_FILE_BYTES
    max_image_bytes: int = MAX_IMAGE_BYTES
    max_audio_bytes: int = MAX_AUDIO_BYTES
    max_per_turn: int = MAX_ATTACHMENTS_PER_TURN


def _safe_name(name: str) -> str:
    # Never use a client-provided path.  Keep Unicode letters while removing
    # control characters and path separators.
    raw = Path(str(name or "attachment")).name
    raw = re.sub(r"[\x00-\x1f\x7f]", "", raw).replace("..", "_")
    return raw[:180] or "attachment"


class AttachmentManager:
    def __init__(self, runtime_root: Path, store: AgentStore, limits: AttachmentLimits | None = None):
        self.runtime_root = Path(runtime_root).resolve()
        self.root = (self.runtime_root / "chat" / "attachments").resolve()
        self.root.mkdir(parents=True, exist_ok=True)
        self.store = store
        self.limits = limits or AttachmentLimits()

    def _ensure_inside(self, path: Path) -> Path:
        resolved = path.resolve()
        try:
            resolved.relative_to(self.root)
        except ValueError as exc:
            raise AttachmentError("附件路径越过运行时附件目录。") from exc
        return resolved

    @staticmethod
    def _mime(filename: str, content_type: str | None) -> str:
        guessed = (content_type or "").split(";", 1)[0].strip().lower()
        return guessed if guessed and guessed != "application/octet-stream" else (mimetypes.guess_type(filename)[0] or "application/octet-stream")

    def validate(self, filename: str, size: int, content_type: str | None = None) -> tuple[str, str]:
        safe = _safe_name(filename)
        suffix = Path(safe).suffix.lower()
        if suffix not in ALLOWED_EXTENSIONS:
            raise AttachmentError(f"不支持的附件类型：{suffix or '无扩展名'}。")
        mime = self._mime(safe, content_type)
        if size <= 0:
            raise AttachmentError("附件不能为空。")
        limit = self.limits.max_file_bytes
        if suffix in IMAGE_EXTENSIONS:
            limit = self.limits.max_image_bytes
        elif suffix in AUDIO_EXTENSIONS:
            limit = self.limits.max_audio_bytes
        if size > limit:
            raise AttachmentError(f"附件超过大小限制（最多 {limit // (1024 * 1024)} MB）。")
        return safe, mime

    def _read_stream(self, stream: BinaryIO, expected_size: int | None = None) -> bytes:
        data = stream.read(self.limits.max_file_bytes + 1)
        if len(data) > self.limits.max_file_bytes:
            raise AttachmentError("附件超过通用大小限制。")
        if expected_size and expected_size != len(data):
            # A mismatching multipart header is not trusted; use actual bytes.
            expected_size = len(data)
        return data

    def save_bytes(self, filename: str, data: bytes, content_type: str | None = None) -> dict[str, Any]:
        safe_name, mime = self.validate(filename, len(data), content_type)
        digest = sha256(data).hexdigest()
        existing = next((a for a in self._all_attachments() if a.get("sha256") == digest), None)
        if existing:
            return self.public(existing)
        attachment_id = self.store.new_id("attachment")
        directory = self._ensure_inside(self.root / digest[:2])
        directory.mkdir(parents=True, exist_ok=True)
        path = self._ensure_inside(directory / f"{attachment_id}_{safe_name}")
        path.write_bytes(data)
        record = self.store.create_attachment({
            "attachment_id": attachment_id, "sha256": digest, "original_name": safe_name,
            "mime_type": mime, "size_bytes": len(data), "storage_path": str(path),
            "parse_status": "pending",
        })
        try:
            text, metadata = self.extract(path, mime)
            self.store.update_attachment_parse(attachment_id, "parsed", text, metadata)
            record = self.store.get_attachment(attachment_id) or record
        except Exception as exc:  # parsing failure must not discard a valid file
            self.store.update_attachment_parse(attachment_id, "failed", "", {"error": str(exc)[:500]})
            record = self.store.get_attachment(attachment_id) or record
        return self.public(record)

    def save_stream(self, filename: str, stream: BinaryIO, content_type: str | None = None, expected_size: int | None = None) -> dict[str, Any]:
        return self.save_bytes(filename, self._read_stream(stream, expected_size), content_type)

    def extract(self, path: Path, mime: str) -> tuple[str, dict[str, Any]]:
        suffix = path.suffix.lower()
        if suffix in IMAGE_EXTENSIONS:
            metadata: dict[str, Any] = {"kind": "image"}
            try:
                from PIL import Image
                with Image.open(path) as image:
                    image.verify()
                with Image.open(path) as image:
                    metadata.update({"width": image.width, "height": image.height, "format": image.format})
            except Exception as exc:
                raise AttachmentError("图片校验失败，文件可能已损坏。") from exc
            return "", metadata
        if suffix in AUDIO_EXTENSIONS:
            metadata = {"kind": "audio", "transcription_status": "not_requested", "duration_verified": False}
            try:
                from mutagen import File as MutagenFile
                media = MutagenFile(str(path))
                duration = float(media.info.length) if media is not None and getattr(media, "info", None) else 0.0
                if duration > 30 * 60:
                    raise AttachmentError("语音附件超过 30 分钟限制。")
                if duration > 0:
                    metadata.update({"duration_seconds": round(duration, 3), "duration_verified": True})
            except AttachmentError:
                raise
            except Exception:
                metadata["duration_warning"] = "无法在本地验证时长，转写前仍会受 Provider 限制。"
            return "", metadata
        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            normalized = text.strip()
            return text[:MAX_EXTRACTED_CHARS], {
                "kind": "document",
                "pages": len(reader.pages),
                "vision_required": not bool(normalized),
                "text_extraction": "available" if normalized else "empty_or_scanned",
            }
        if suffix == ".docx":
            from docx import Document
            document = Document(str(path))
            return "\n".join(p.text for p in document.paragraphs)[:MAX_EXTRACTED_CHARS], {"kind": "document"}
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            workbook = load_workbook(str(path), read_only=True, data_only=True)
            chunks = []
            for sheet in workbook.worksheets:
                chunks.append(f"[{sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    chunks.append("\t".join("" if value is None else str(value) for value in row))
            return "\n".join(chunks)[:MAX_EXTRACTED_CHARS], {"kind": "spreadsheet", "sheets": workbook.sheetnames}
        if suffix == ".pptx":
            try:
                from pptx import Presentation
            except ImportError as exc:
                raise AttachmentError("PPTX 提取需要安装 python-pptx。") from exc
            presentation = Presentation(str(path))
            chunks = []
            for index, slide in enumerate(presentation.slides, 1):
                chunks.append(f"[slide {index}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        chunks.append(str(shape.text))
            return "\n".join(chunks)[:MAX_EXTRACTED_CHARS], {"kind": "presentation", "slides": len(presentation.slides)}
        raw = path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".json":
            try:
                raw = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                pass
        return raw[:MAX_EXTRACTED_CHARS], {"kind": "text"}

    def get(self, attachment_id: str, include_text: bool = False) -> dict[str, Any] | None:
        record = self.store.get_attachment(attachment_id)
        if not record:
            return None
        return self.public(record, include_text=include_text)

    def delete(self, attachment_id: str) -> dict[str, Any] | None:
        record = self.store.delete_attachment(attachment_id)
        if not record:
            return None
        path = Path(record["storage_path"])
        try:
            self._ensure_inside(path).unlink(missing_ok=True)
            if path.parent != self.root and not any(path.parent.iterdir()):
                path.parent.rmdir()
        except (OSError, AttachmentError):
            pass
        return self.public(record)

    def set_retention(self, attachment_id: str, days: int | None) -> dict[str, Any] | None:
        if days is not None and not 1 <= int(days) <= 3650:
            raise AttachmentError("附件保留期限必须在 1 到 3650 天之间，或设为永久。")
        expires_at = None
        if days is not None:
            from datetime import timedelta
            expires_at = (datetime.now(timezone.utc) + timedelta(days=int(days))).isoformat()
        record = self.store.update_attachment_expiry(attachment_id, expires_at)
        return self.public(record) if record else None

    def cleanup_expired(self) -> dict[str, int]:
        now = datetime.now(timezone.utc)
        deleted = 0
        for record in self.store.list_attachments(limit=500):
            raw = str(record.get("expires_at") or "")
            if not raw:
                continue
            try:
                expires = datetime.fromisoformat(raw.replace("Z", "+00:00"))
            except ValueError:
                continue
            if expires <= now and self.delete(str(record["attachment_id"])):
                deleted += 1
        return {"deleted": deleted}

    def _all_attachments(self) -> list[dict[str, Any]]:
        # The store intentionally exposes point lookups only.  For deduplication
        # this bounded query avoids leaking metadata through the API.
        with self.store._connect() as connection:  # internal, same process/store
            rows = connection.execute("SELECT * FROM attachments").fetchall()
        return [self.store.attachment_row(row) for row in rows]

    @staticmethod
    def public(record: dict[str, Any], include_text: bool = False) -> dict[str, Any]:
        result = {key: record.get(key) for key in ("attachment_id", "sha256", "original_name", "mime_type", "size_bytes", "parse_status", "metadata", "created_at", "updated_at", "expires_at")}
        if include_text:
            result["extracted_text"] = record.get("extracted_text", "")
        return result
